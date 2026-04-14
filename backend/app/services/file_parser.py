"""
app/services/file_parser.py
───────────────────────────
Handles extraction of raw text from all supported file types.

Supported:
  .txt   → plain read
  .pdf   → PyMuPDF (fitz)
  .docx  → python-docx
  .xlsx  → pandas + openpyxl
  .pptx  → python-pptx
"""

import io
import pandas as pd
from pathlib import Path
from loguru import logger


# ── PDF ───────────────────────────────────────────────────────────────────────

def parse_pdf(content: bytes) -> tuple[str, int]:
    """Extract text from PDF bytes. Returns (text, page_count)."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=content, filetype="pdf")
        pages: list[str] = []
        for page_num, page in enumerate(doc):
            text = page.get_text("text")
            if text.strip():
                pages.append(f"[Page {page_num + 1}]\n{text.strip()}")
        doc.close()
        return "\n\n".join(pages), len(doc)
    except Exception as exc:
        logger.error(f"PDF parsing failed: {exc}")
        raise ValueError(f"Could not parse PDF: {exc}") from exc


# ── DOCX ──────────────────────────────────────────────────────────────────────

def parse_docx(content: bytes) -> tuple[str, None]:
    """Extract text from .docx bytes. Returns (text, None)."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        paragraphs: list[str] = []

        # Body paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                # Preserve heading level hints
                if para.style.name.startswith("Heading"):
                    paragraphs.append(f"\n## {para.text.strip()}\n")
                else:
                    paragraphs.append(para.text.strip())

        # Tables
        for table in doc.tables:
            rows: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                paragraphs.append("\n[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]")

        return "\n\n".join(paragraphs), None
    except Exception as exc:
        logger.error(f"DOCX parsing failed: {exc}")
        raise ValueError(f"Could not parse DOCX: {exc}") from exc


# ── XLSX ──────────────────────────────────────────────────────────────────────

def parse_xlsx(content: bytes) -> tuple[str, None]:
    """Extract tabular data from .xlsx bytes. Returns (text, None)."""
    try:
        xl = pd.ExcelFile(io.BytesIO(content), engine="openpyxl")
        sheet_texts: list[str] = []

        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name, dtype=str)
            df.fillna("", inplace=True)

            if df.empty:
                continue

            # Convert to readable markdown-style table
            header = " | ".join(str(c) for c in df.columns)
            separator = " | ".join(["---"] * len(df.columns))
            rows = [" | ".join(str(v) for v in row) for _, row in df.iterrows()]
            table_text = f"### Sheet: {sheet_name}\n{header}\n{separator}\n" + "\n".join(rows)
            sheet_texts.append(table_text)

        return "\n\n".join(sheet_texts), None
    except Exception as exc:
        logger.error(f"XLSX parsing failed: {exc}")
        raise ValueError(f"Could not parse XLSX: {exc}") from exc


# ── PPTX ──────────────────────────────────────────────────────────────────────

def parse_pptx(content: bytes) -> tuple[str, int]:
    """Extract text from .pptx bytes. Returns (text, slide_count)."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides_text: list[str] = []

        for idx, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if slide_parts:
                slides_text.append(f"[Slide {idx}]\n" + "\n".join(slide_parts))

        return "\n\n".join(slides_text), len(prs.slides)
    except Exception as exc:
        logger.error(f"PPTX parsing failed: {exc}")
        raise ValueError(f"Could not parse PPTX: {exc}") from exc


# ── TXT ───────────────────────────────────────────────────────────────────────

def parse_txt(content: bytes) -> tuple[str, None]:
    """Decode plain text bytes. Tries UTF-8 then Latin-1 fallback."""
    try:
        try:
            return content.decode("utf-8"), None
        except UnicodeDecodeError:
            return content.decode("latin-1"), None
    except Exception as exc:
        raise ValueError(f"Could not decode text file: {exc}") from exc


# ── Dispatcher ────────────────────────────────────────────────────────────────

def parse_document(filename: str, content: bytes) -> tuple[str, int | None]:
    """
    Route file to correct parser based on extension.

    Returns:
        (extracted_text, page_or_slide_count | None)
    """
    ext = Path(filename).suffix.lower()
    logger.info(f"Parsing '{filename}' as {ext} ({len(content):,} bytes)")

    parsers = {
        ".pdf":  parse_pdf,
        ".docx": parse_docx,
        ".xlsx": parse_xlsx,
        ".pptx": parse_pptx,
        ".txt":  parse_txt,
    }

    if ext not in parsers:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Supported: {', '.join(parsers.keys())}"
        )

    text, count = parsers[ext](content)

    if not text.strip():
        raise ValueError(
            "Document appears to be empty or contains no extractable text. "
            "Scanned image-only PDFs are not supported."
        )

    logger.success(
        f"Parsed '{filename}': {len(text):,} chars extracted"
        + (f", {count} pages/slides" if count else "")
    )
    return text, count
